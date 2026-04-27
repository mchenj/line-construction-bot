"""
weekly_phase1.py
Phase 1: เติมข้อมูลรายงานประจำสัปดาห์จาก daily reports
- 04 รายละเอียดโครงการ.docx → fill 3 tables (บุคลากร, บันทึก, อากาศ)
- ภาคผนวก 1 ภาพถ่าย.docx → fill photo grid
- ภาคผนวก 3 รายงานประจำวัน → loop generate_daily 8 ใบ
ส่งกลับเป็น ZIP file รวมทุก output
"""

import io, os, re, json, zipfile
from copy import deepcopy
from datetime import date, timedelta
from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

# template paths
_TEMPLATE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "templates_weekly")
TEMPLATE_COVER          = os.path.join(_TEMPLATE_DIR, "01_cover.pptx")
TEMPLATE_TOC            = os.path.join(_TEMPLATE_DIR, "03_toc.docx")
TEMPLATE_PROJECT_DETAILS= os.path.join(_TEMPLATE_DIR, "04_project_details.docx")
TEMPLATE_APPENDIX1_PHOTOS = os.path.join(_TEMPLATE_DIR, "appendix1_photos.docx")

THAI_MONTHS_FULL = ["","มกราคม","กุมภาพันธ์","มีนาคม","เมษายน","พฤษภาคม","มิถุนายน",
                    "กรกฎาคม","สิงหาคม","กันยายน","ตุลาคม","พฤศจิกายน","ธันวาคม"]
THAI_MONTHS_ABBR = ["","ม.ค.","ก.พ.","มี.ค.","เม.ย.","พ.ค.","มิ.ย.",
                    "ก.ค.","ส.ค.","ก.ย.","ต.ค.","พ.ย.","ธ.ค."]
THAI_DIGITS = "๐๑๒๓๔๕๖๗๘๙"


def _to_thai_digits(s) -> str:
    return "".join(THAI_DIGITS[int(c)] if c.isdigit() else c for c in str(s))


def _thai_date_full(d: date) -> str:
    """24 มีนาคม 2569"""
    return f"{d.day} {THAI_MONTHS_FULL[d.month]} {d.year + 543}"


def _thai_date_full_thai_digits(d: date) -> str:
    """๒๔ มีนาคม ๒๕๖๙"""
    return f"{_to_thai_digits(d.day)} {THAI_MONTHS_FULL[d.month]} {_to_thai_digits(d.year + 543)}"


def _thai_date_short(d: date) -> str:
    """24 มี.ค. 69"""
    return f"{d.day} {THAI_MONTHS_ABBR[d.month]} {str(d.year + 543)[2:]}"


# ════════════════════════════════════════
# Cell helpers
# ════════════════════════════════════════

def _set_page_break_before(para):
    """เพิ่ม <w:pageBreakBefore/> ใน pPr → หัวข้อนี้ขึ้นหน้าใหม่เสมอ"""
    pPr = para._p.get_or_add_pPr()
    # ลบ pageBreakBefore เก่าก่อนถ้ามี
    for pb in pPr.findall(qn("w:pageBreakBefore")):
        pPr.remove(pb)
    pb = OxmlElement("w:pageBreakBefore")
    pPr.insert(0, pb)


def force_page_break_before_heading(doc, heading_keyword: str, exclude_keyword: str = None):
    """หา paragraph ที่มี heading_keyword (และไม่มี exclude_keyword) แล้วใส่ page break"""
    for p in doc.paragraphs:
        text = p.text.strip()
        if heading_keyword in text:
            if exclude_keyword and exclude_keyword in text:
                continue
            _set_page_break_before(p)
            return True
    return False


def _set_cell_text(cell, text: str, font_size: int = 14, bold: bool = False, center: bool = True):
    """ล้าง cell แล้วใส่ข้อความใหม่ พร้อมจัด format"""
    # ลบ paragraph เก่าทั้งหมดยกเว้นอันแรก
    p = cell.paragraphs[0]
    for old_p in cell.paragraphs[1:]:
        old_p._element.getparent().remove(old_p._element)
    # ล้าง runs ใน paragraph แรก
    for r in p.runs:
        r._element.getparent().remove(r._element)
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER if center else WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(text)
    run.font.name = "TH SarabunIT๙"
    run.font.size = Pt(font_size)
    run.bold = bold
    # set CS font (สำหรับภาษาไทย)
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.append(rFonts)
    rFonts.set(qn("w:cs"), "TH SarabunIT๙")
    rFonts.set(qn("w:hAnsi"), "TH SarabunIT๙")
    rFonts.set(qn("w:ascii"), "TH SarabunIT๙")


# ════════════════════════════════════════
# Table 2: บุคลากรและเครื่องจักร 8 วัน
# ════════════════════════════════════════

# mapping: ชื่อใน DB/equipment list → row label ในตาราง
_PERSONNEL_EQUIP_ROWS = [
    # (label, get_func)  — get_func รับ daily_data และคืน int
    ("วิศวกร/ช่าง",          lambda d: (d.get("engineers") or 0) + (d.get("skilled_workers") or 0)),
    ("หัวหน้าคนงาน",          lambda d: d.get("foremen") or 0),
    ("กรรมกร",                lambda d: d.get("laborers") or 0),
    ("รถเฮี้ยบ",              lambda d: _equip_qty(d, "รถเฮี้ยบ")),
    ("รถเทเลอร์",             lambda d: _equip_qty(d, "รถเทเลอร์")),
    ("รถสกัดคอนกรีต",         lambda d: _equip_qty(d, "รถสกัดคอนกรีต")),
    ("รถเกรด",                lambda d: _equip_qty(d, "รถเกรด")),
    ("รถแบ็คโฮ",              lambda d: _equip_qty(d, "รถแบ็คโฮ")),
    ("รถบด",                  lambda d: _equip_qty(d, "รถบด")),
    ("รถน้ำ",                 lambda d: _equip_qty(d, "รถน้ำ")),
    ("รถบรรทุก",              lambda d: _equip_qty(d, "รถบรรทุก")),
    ("รถแทร็กเตอร์",          lambda d: _equip_qty(d, "รถแทร็กเตอร์")),
    ("รถเครน",                lambda d: _equip_qty(d, "รถเครน")),
    ("รถบดล้อยาง",            lambda d: _equip_qty(d, "รถบดล้อยาง")),
    ("ปั่นจั่น",              lambda d: _equip_qty(d, "ปั่นจั่น")),
    ("รถสกัดคอนกรีตเสาเข็ม",  lambda d: _equip_qty(d, "รถสกัดคอนกรีตเสาเข็ม")),
    ("รถบรรทุก 10 ล้อ",       lambda d: _equip_qty(d, "รถบรรทุก 10 ล้อ")),
    ("รถบรรทุก 6 ล้อ",        lambda d: _equip_qty(d, "รถบรรทุก 6 ล้อ")),
    ("กล้องสำรวจแนว",         lambda d: _equip_qty(d, "กล้องสำรวจแนว")),
    ("กล้องระดับ",            lambda d: _equip_qty(d, "กล้องระดับ")),
    ("เครื่องจี้คอนกรีต",     lambda d: _equip_qty(d, "เครื่องจี้คอนกรีต")),
    ("เครื่องเชื่อม",         lambda d: _equip_qty(d, "เครื่องเชื่อม")),
    ("เครื่องตบดิน",          lambda d: _equip_qty(d, "เครื่องตบดิน")),
    ("เครื่องสูบน้ำ",         lambda d: _equip_qty(d, "เครื่องสูบน้ำ")),
    ("รถบดสั่นสะเทือน",       lambda d: _equip_qty(d, "รถบดสั่นสะเทือน")),
    ("รถน้ำ",                 lambda d: _equip_qty(d, "รถน้ำ")),  # ซ้ำตามรูปแบบ template
    ("รถPRIME COAT",          lambda d: _equip_qty(d, "รถPRIME COAT")),
    ("รถPAVE ยาง",            lambda d: _equip_qty(d, "รถPAVE ยาง")),
]


def _parse_equipment(daily_data: dict) -> dict:
    """แปลง equipment field (str/list) → dict {name: qty}"""
    raw = daily_data.get("equipment")
    if isinstance(raw, str):
        try:
            raw = json.loads(raw)
        except Exception:
            raw = []
    if not raw:
        return {}
    return {(e.get("name") or "").strip().replace("นํ้า", "น้ำ"): int(e.get("qty") or 0)
            for e in raw}


def _equip_qty(daily_data: dict, name: str) -> int:
    """ดึง qty ของเครื่องจักร 'name' จาก daily_data"""
    eq = _parse_equipment(daily_data)
    name_norm = name.replace("นํ้า", "น้ำ")
    # exact match ก่อน
    if name_norm in eq:
        return eq[name_norm]
    # partial match (e.g. "รถบรรทุก" จะ match "รถบรรทุก 10 ล้อ" ด้วย — ต้องระวัง)
    # ใน template มีทั้ง "รถบรรทุก" และ "รถบรรทุก 10 ล้อ" แยกกัน
    # ดังนั้นต้อง exact match เท่านั้น
    return 0


def fill_personnel_equipment_table(table, daily_list: list):
    """เติม Table 2 (บุคลากรและเครื่องจักร 8 วัน) — 31 rows × 10 cols
    rows[0-1]: header (skip)
    rows[2..29]: 28 รายการ (3 personnel + 25 equipment)
    rows[30]: รวม
    daily_list ต้องมี 8 รายการ เรียงจากวันแรกถึงวันสุดท้าย
    """
    # row 0: header (merged)
    # row 1: ลำดับที่ | ประเภทบุคลากร | 16 | 17 | ...
    # row 2 onwards: data
    n_days = len(daily_list)

    # update header row 1 — set day numbers (cols 2..9)
    if len(table.rows) > 1:
        hdr = table.rows[1]
        for i in range(min(n_days, 8)):
            d = date.fromisoformat(daily_list[i].get("work_date"))
            if 2 + i < len(hdr.cells):
                _set_cell_text(hdr.cells[2 + i], str(d.day), font_size=14, bold=True)

    # data rows: index 2 onwards
    n_rows_avail = len(table.rows) - 3  # หัก header 2 rows + รวม 1 row
    n_items = min(len(_PERSONNEL_EQUIP_ROWS), n_rows_avail)

    daily_totals = [0] * 8

    for ri in range(n_items):
        label, get_fn = _PERSONNEL_EQUIP_ROWS[ri]
        row = table.rows[2 + ri]
        # col 0: ลำดับที่
        _set_cell_text(row.cells[0], str(ri + 1), font_size=14)
        # col 1: ประเภท
        _set_cell_text(row.cells[1], label, font_size=14, center=False)
        # cols 2..9: qty per day
        for di in range(min(n_days, 8)):
            qty = get_fn(daily_list[di])
            daily_totals[di] += qty
            if 2 + di < len(row.cells):
                _set_cell_text(row.cells[2 + di], str(qty) if qty else "0", font_size=14)

    # รวม row (last row)
    if len(table.rows) >= 3 + n_items:
        sum_row_idx = 2 + n_items
        # หา "รวม" row — ในตัวอย่างคือ row 30
        sum_row = table.rows[-1]
        # col 0-1 may be merged → แค่ใส่ "รวม"
        _set_cell_text(sum_row.cells[0], "", font_size=14)
        _set_cell_text(sum_row.cells[1], "รวม", font_size=14, bold=True)
        for di in range(min(n_days, 8)):
            if 2 + di < len(sum_row.cells):
                _set_cell_text(sum_row.cells[2 + di], str(daily_totals[di]),
                               font_size=14, bold=True)


# ════════════════════════════════════════
# Table 3: บันทึกการปฏิบัติงาน
# ════════════════════════════════════════

def fill_diary_table(table, daily_list: list):
    """เติม Table 3 (บันทึกการปฏิบัติงาน) — 28 rows × 3 cols
    Format:
      วัน เดือน ปี | รายละเอียดงาน | สถานที่
      ๑๖ มีนาคม ๒๕๖๙ | (activity 1) |
                     | (activity 2) |
                     | (activity 3) |
      ๑๗ มีนาคม ๒๕๖๙ | (activity 1) |
      ...
    """
    # row 0 = header (skip)
    # ลบ data rows ทั้งหมด เก็บแค่ header
    template_row = None
    if len(table.rows) > 1:
        template_row = deepcopy(table.rows[1]._tr)
    # ลบ rows ตั้งแต่ index 1 เป็นต้นไป
    while len(table.rows) > 1:
        last_row = table.rows[-1]
        last_row._element.getparent().remove(last_row._element)

    if template_row is None:
        return

    # สร้าง rows ใหม่ตาม daily_list
    for daily in daily_list:
        d = date.fromisoformat(daily.get("work_date"))
        date_str = _thai_date_full_thai_digits(d)
        activities = daily.get("activities") or []

        if not activities:
            # อย่างน้อยมี 1 row ต่อวัน เพื่อใส่วันที่
            new_tr = deepcopy(template_row)
            table._tbl.append(new_tr)
            new_row = table.rows[-1]
            _set_cell_text(new_row.cells[0], date_str, font_size=14, center=False)
            _set_cell_text(new_row.cells[1], "—", font_size=14, center=False)
            _set_cell_text(new_row.cells[2], "", font_size=14, center=False)
            continue

        for ai, act in enumerate(activities):
            new_tr = deepcopy(template_row)
            table._tbl.append(new_tr)
            new_row = table.rows[-1]
            # col 0: วัน เดือน ปี (เฉพาะ row แรกของวัน)
            _set_cell_text(new_row.cells[0], date_str if ai == 0 else "",
                           font_size=14, center=False)
            # col 1: รายละเอียดงาน
            desc = act.get("desc") or act.get("description") or ""
            _set_cell_text(new_row.cells[1], desc, font_size=14, center=False)
            # col 2: สถานที่ (ว่างไว้)
            _set_cell_text(new_row.cells[2], "", font_size=14, center=False)


# ════════════════════════════════════════
# Table 4: สภาพอากาศและระดับน้ำ
# ════════════════════════════════════════

_WEATHER_LABEL_MAP = {
    "แจ่มใส": "แจ่มใส",
    "เมฆมาก": "มืดครึ้ม",
    "มืดครึ้ม": "มืดครึ้ม",
    "ฝนตกเล็กน้อย": "ฝนตกเล็กน้อย",
    "ฝนตกหนัก": "ฝนตกหนัก",
}


def _short_weather(weather: str) -> str:
    if not weather:
        return "—"
    for k, v in _WEATHER_LABEL_MAP.items():
        if k in weather:
            return v
    return weather[:10]


def fill_weather_table(table, daily_list: list):
    """เติม Table 4 (สภาพอากาศและระดับน้ำ) — 6 rows × 10 cols
    row[0]: วันที่ | 16 | 17 | 18 | 19 | 20 | 21 | 22 | 23
    row[1-3]: สภาพอากาศ (3 rows ในตัวอย่าง — ใส่บรรทัดเดียวก็ได้)
    row[4]: ระดับน้ำสูงสุด | +92.60 x 8
    row[5]: หมายเหตุ
    """
    n_days = len(daily_list)

    # row 0: update day numbers (cols 1..8)
    if len(table.rows) > 0:
        hdr = table.rows[0]
        for i in range(min(n_days, 8)):
            d = date.fromisoformat(daily_list[i].get("work_date"))
            if 1 + i < len(hdr.cells):
                _set_cell_text(hdr.cells[1 + i], str(d.day), font_size=14, bold=True)

    # row 1: สภาพอากาศ — note: rows 1-3 ถูก merge แนวตั้งใน template
    # แค่ใส่ที่ row[1] ก็พอ ไม่ต้อง clear row 2,3 (จะลบข้อมูลใน merged cell)
    if len(table.rows) > 1:
        wrow = table.rows[1]
        for i in range(min(n_days, 8)):
            wx = _short_weather(daily_list[i].get("weather_morning") or "")
            if 1 + i < len(wrow.cells):
                _set_cell_text(wrow.cells[1 + i], wx, font_size=13)

    # row 4: ระดับน้ำสูงสุด
    if len(table.rows) > 4:
        lrow = table.rows[4]
        for i in range(min(n_days, 8)):
            wl = daily_list[i].get("water_level")
            wl_str = (f"+{wl:.2f}" if wl is not None and wl >= 0
                      else (f"{wl:.2f}" if wl is not None else "—"))
            if 1 + i < len(lrow.cells):
                _set_cell_text(lrow.cells[1 + i], wl_str, font_size=13)


# ════════════════════════════════════════
# ภาคผนวก 1: ภาพถ่ายการปฏิบัติงาน
# ════════════════════════════════════════

async def _download_image(url: str) -> bytes:
    import httpx
    if not url:
        return b""
    try:
        async with httpx.AsyncClient(timeout=20.0) as client:
            r = await client.get(url)
            if r.status_code == 200:
                return r.content
    except Exception as e:
        print(f"download image failed: {e}")
    return b""


def _add_image_caption(doc_or_body, caption: str, font_size: int = 16):
    """เพิ่ม caption ใต้ภาพ — สไตล์เดียวกับรายงานประจำวัน
    (ขึ้นบรรทัดใหม่ใช้ line break ใน paragraph เดียว ไม่ใช่ paragraph ใหม่ → ลดช่องว่าง)
    """
    cap_p = doc_or_body.add_paragraph()
    cap_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap_p.paragraph_format.space_before = Pt(4)
    cap_p.paragraph_format.space_after = Pt(6)
    lines = [l for l in caption.split('\n') if l.strip()]
    for i, line in enumerate(lines):
        run = cap_p.add_run(line)
        run.font.name = "TH SarabunIT๙"
        run.font.size = Pt(font_size)
        rPr = run._r.get_or_add_rPr()
        rFonts = rPr.find(qn("w:rFonts"))
        if rFonts is None:
            rFonts = OxmlElement("w:rFonts")
            rPr.append(rFonts)
        for attr in ("ascii", "hAnsi", "cs"):
            rFonts.set(qn(f"w:{attr}"), "TH SarabunIT๙")
        szCs = rPr.find(qn("w:szCs"))
        if szCs is None:
            szCs = OxmlElement("w:szCs")
            rPr.append(szCs)
        szCs.set(qn("w:val"), str(font_size * 2))
        if i < len(lines) - 1:
            run.add_break()


async def fill_photos_table(doc, daily_list: list, week_no: int, week_start: date, week_end: date):
    """เติมภาพถ่ายในภาคผนวก 1 — สไตล์เดียวกับรายงานประจำวัน
    template:
      paragraph[0]: ภาพถ่ายการปฏิบัติงานประจำสัปดาห์ที่ XX/YYYY (วันที่ DD - DD MMMM YYYY)
      table[0]: ตารางภาพ (จะถูกลบทิ้ง แล้ว append รูปแบบ full-width แทน)
    """
    # update title paragraph (คงรูปแบบเดิม)
    if doc.paragraphs:
        p = doc.paragraphs[0]
        for r in p.runs:
            r._element.getparent().remove(r._element)
        title = (f"ภาพถ่ายการปฏิบัติงานประจำสัปดาห์ที่ {week_no}/{week_start.year + 543}"
                 f" (วันที่ {week_start.day} – {week_end.day} {THAI_MONTHS_FULL[week_start.month]}"
                 f" {week_start.year + 543})")
        run = p.add_run(title)
        run.bold = True
        run.font.name = "TH SarabunIT๙"
        run.font.size = Pt(16)

    # ลบ table เดิมใน template (เปลี่ยนเป็น layout แบบรูปต่อ paragraph)
    for tbl in list(doc.tables):
        tbl._element.getparent().remove(tbl._element)

    # collect (img_url, caption, work_date_obj) tuples
    photos = []
    for daily in daily_list:
        try:
            d = date.fromisoformat(daily.get("work_date"))
        except Exception:
            continue
        for img_info in (daily.get("images") or []):
            url = img_info.get("url") or img_info.get("image_url")
            cap = (img_info.get("caption") or "").strip()
            if url:
                photos.append((url, cap, d))

    # เติมรูปทีละใบ (full width 5 นิ้ว) + caption — แบบเดียวกับ daily
    for url, cap, d in photos:
        img_bytes = await _download_image(url)
        if not img_bytes:
            continue
        img_p = doc.add_paragraph()
        img_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        try:
            img_p.add_run().add_picture(io.BytesIO(img_bytes), width=Inches(5.0))
        except Exception as e:
            img_p.add_run(f"[image error: {e}]")
            continue
        date_th = _thai_date_full(d)
        caption_text = f"วันที่ {date_th}"
        if cap:
            caption_text = f"{cap}\nวันที่ {date_th}"
        _add_image_caption(doc, caption_text, font_size=16)


# ════════════════════════════════════════
# หน้าปก (PPTX)
# ════════════════════════════════════════

def _set_pptx_run_text(run, new_text: str):
    """แทนข้อความใน run ของ pptx โดยรักษา format เดิม"""
    run.text = new_text


def fill_cover(prs, week_no: int, week_start: date, week_end: date):
    """เติมหน้าปก PPTX:
    shape[3].para[0]: รายงานประจำสัปดาห์ที่ XX/YYYY
    shape[3].para[1]: (ประจำวันที่ DD – DD เดือน YYYY)
    """
    if not prs.slides:
        return
    slide = prs.slides[0]
    year_th = week_start.year + 543
    new_title = f"รายงานประจำสัปดาห์ที่ {week_no}/{year_th}"
    new_subtitle = (f"(ประจำวันที่ {week_start.day} – {week_end.day} "
                    f"{THAI_MONTHS_FULL[week_end.month]} {year_th})")

    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        for p in tf.paragraphs:
            full_text = "".join(r.text for r in p.runs)
            if "รายงานประจำสัปดาห์ที่" in full_text and "/" in full_text:
                # เคลียร์ runs เก่าและใส่ใหม่ใน run แรก
                if p.runs:
                    p.runs[0].text = new_title
                    for r in p.runs[1:]:
                        r.text = ""
            elif "ประจำวันที่" in full_text:
                if p.runs:
                    p.runs[0].text = new_subtitle
                    for r in p.runs[1:]:
                        r.text = ""


# ════════════════════════════════════════
# สารบัญ (DOCX)
# ════════════════════════════════════════

def _replace_para_text_keep_format(para, new_text: str):
    """แทนข้อความใน paragraph โดย keep format ของ run แรก ลบ runs อื่นทิ้ง"""
    if not para.runs:
        return
    para.runs[0].text = new_text
    # ลบ runs ที่เหลือ
    for r in para.runs[1:]:
        r._element.getparent().remove(r._element)


def fill_toc(doc, week_no: int, year_th: int):
    """เติมสารบัญ:
    paragraph[0]: รายงานความก้าวหน้าประจำสัปดาห์ที่ XX/YYYY
    """
    new_title = f"รายงานความก้าวหน้าประจำสัปดาห์ที่ {week_no}/{year_th}"
    for p in doc.paragraphs[:5]:
        if "รายงานความก้าวหน้าประจำสัปดาห์ที่" in p.text:
            _replace_para_text_keep_format(p, new_title)
            return


# ════════════════════════════════════════
# Main: generate Phase 1 weekly report
# ════════════════════════════════════════

def fill_template_captions(doc, week_no: int, ws_date: date, we_date: date):
    """แก้ caption "สัปดาห์ที่ XX/YYYY ระหว่างวันที่ DD - DD เดือน YYYY" และ
    "ประจำวันที่ DD - DD เดือน YYYY" ใน 04_project_details.docx
    """
    year_th = ws_date.year + 543
    # ถ้าข้ามเดือน ใช้ "16 มี.ค. - 23 เม.ย." style
    if ws_date.month == we_date.month:
        date_str = f"{ws_date.day} – {we_date.day} {THAI_MONTHS_FULL[we_date.month]} {year_th}"
    else:
        date_str = (f"{ws_date.day} {THAI_MONTHS_FULL[ws_date.month]} – "
                    f"{we_date.day} {THAI_MONTHS_FULL[we_date.month]} {year_th}")
    new_full = f"สัปดาห์ที่ {week_no}/{year_th} ระหว่างวันที่ {date_str}"
    new_short = f"ผลการดำเนินงานก่อสร้างประจำวันที่ {date_str}"

    # 1) update paragraphs in body
    for p in doc.paragraphs:
        text = p.text
        if "สัปดาห์ที่" in text and "ระหว่างวันที่" in text:
            _replace_para_text_keep_format(p, new_full)
        elif "ผลการดำเนินงานก่อสร้างประจำวันที่" in text:
            _replace_para_text_keep_format(p, new_short)

    # 2) update Table[2] R0 caption (merged across cols) — ภาคผนวก 4 personnel header
    if len(doc.tables) > 2:
        tbl2 = doc.tables[2]
        if len(tbl2.rows) > 0:
            r0 = tbl2.rows[0]
            for cell in r0.cells:
                if "สัปดาห์ที่" in cell.text and "ระหว่างวันที่" in cell.text:
                    # rebuild cell content
                    p = cell.paragraphs[0]
                    for old_p in cell.paragraphs[1:]:
                        old_p._element.getparent().remove(old_p._element)
                    if p.runs:
                        p.runs[0].text = new_full
                        for r in p.runs[1:]:
                            r._element.getparent().remove(r._element)
                    break  # merged cells share content; update once is enough


async def generate_weekly_phase1(week_no: int, week_start: str, daily_list: list,
                                 project_name: str = "โครงการก่อสร้าง",
                                 week_end: str = None) -> bytes:
    """
    Generate Phase 1 weekly report files และ pack เป็น ZIP

    Args:
        week_no: เลขสัปดาห์ที่ (เช่น 85)
        week_start: วันแรกของสัปดาห์ "YYYY-MM-DD" (เช่น "2026-03-16")
        daily_list: list ของ daily_data dicts (ควรมี 8 วัน)
        project_name: ชื่อโครงการ

    Returns:
        bytes ของไฟล์ ZIP ที่บรรจุ:
            - 02_รายละเอียดโครงการ.docx
            - 03_ภาคผนวก_1_ภาพถ่าย.docx
            - 04_ภาคผนวก_3_รายงานประจำวัน_DD.docx (× n daily)
    """
    from report_generator import generate_daily

    ws = date.fromisoformat(week_start)
    if week_end:
        we = date.fromisoformat(week_end)
    elif daily_list:
        we = ws + timedelta(days=len(daily_list) - 1)
    else:
        we = ws + timedelta(days=7)

    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:

        # ───── 0. หน้าปก (PPTX) ─────
        try:
            from pptx import Presentation
            prs = Presentation(TEMPLATE_COVER)
            fill_cover(prs, week_no, ws, we)
            buf = io.BytesIO()
            prs.save(buf)
            zf.writestr(f"00_หน้าปก_W{week_no}.pptx", buf.getvalue())
        except Exception as e:
            zf.writestr(f"ERROR_cover.txt", f"Error: {e}")

        # ───── 1. สารบัญ (DOCX) ─────
        try:
            doc = Document(TEMPLATE_TOC)
            fill_toc(doc, week_no, ws.year + 543)
            buf = io.BytesIO()
            doc.save(buf)
            zf.writestr(f"01_สารบัญ_W{week_no}.docx", buf.getvalue())
        except Exception as e:
            zf.writestr(f"ERROR_toc.txt", f"Error: {e}")

        # ───── 2. รายละเอียดโครงการ + 3 tables ─────
        try:
            doc = Document(TEMPLATE_PROJECT_DETAILS)
            tables = doc.tables
            # Table[2]: บุคลากร, Table[3]: บันทึก, Table[4]: อากาศ
            if len(tables) > 2:
                fill_personnel_equipment_table(tables[2], daily_list)
            if len(tables) > 3:
                fill_diary_table(tables[3], daily_list)
            if len(tables) > 4:
                fill_weather_table(tables[4], daily_list)
            # Phase 3: ผลการดำเนินงาน (Tables 0, 1) จาก data/construction_plan.xlsx
            try:
                from weekly_phase3 import (read_progress_detail, compute_progress_summary,
                                           fill_progress_summary_table, fill_progress_detail_table)
                detail = read_progress_detail()
                if detail:
                    summary = compute_progress_summary(detail)
                    if len(tables) > 0:
                        fill_progress_summary_table(tables[0], summary)
                    if len(tables) > 1:
                        fill_progress_detail_table(tables[1], detail)
            except Exception as e:
                print(f"⚠️ Phase 3 progress fill skipped: {e}")
            # อัปเดต captions ในเทมเพลต (page 9, 10, 14) ให้ตรงกับ week_no + ช่วงวันที่ใหม่
            fill_template_captions(doc, week_no, ws, we)
            # บังคับให้หัวข้อ "บันทึกการปฏิบัติงานผู้รับจ้าง" (ก่อน table 3) ขึ้นหน้าใหม่เสมอ
            force_page_break_before_heading(doc, "บันทึกการปฏิบัติงานผู้รับจ้าง",
                                            exclude_keyword="ของ")
            buf = io.BytesIO()
            doc.save(buf)
            zf.writestr(f"02_รายละเอียดโครงการ_W{week_no}.docx", buf.getvalue())
        except Exception as e:
            zf.writestr(f"ERROR_project_details.txt", f"Error: {e}")

        # ───── 2. ภาคผนวก 1: ภาพถ่าย ─────
        try:
            doc = Document(TEMPLATE_APPENDIX1_PHOTOS)
            await fill_photos_table(doc, daily_list, week_no, ws, we)
            buf = io.BytesIO()
            doc.save(buf)
            zf.writestr(f"03_ภาคผนวก_1_ภาพถ่าย_W{week_no}.docx", buf.getvalue())
        except Exception as e:
            zf.writestr(f"ERROR_appendix1_photos.txt", f"Error: {e}")

        # ───── 3. ภาคผนวก 3: รายงานประจำวัน 8 ใบ ─────
        for i, daily in enumerate(daily_list):
            try:
                wd = daily.get("work_date")
                d = date.fromisoformat(wd)
                fb = await generate_daily(wd, daily, project_name)
                fname = f"04_ภาคผนวก_3_รายงานประจำวัน_{d.strftime('%Y%m%d')}.docx"
                zf.writestr(fname, fb)
            except Exception as e:
                zf.writestr(f"ERROR_daily_{i+1}.txt", f"Error: {e}")

        # ───── 4. ภาคผนวก 4: บุคลากร CM (Phase 3) ─────
        try:
            from weekly_phase3 import read_cm_personnel, fill_appendix4_xlsx, TEMPLATE_APPENDIX4_XLSX
            cm_data = read_cm_personnel(ws, we)
            if cm_data.get("personnel"):
                fb = fill_appendix4_xlsx(TEMPLATE_APPENDIX4_XLSX, week_no, ws, we, cm_data)
                zf.writestr(f"05_ภาคผนวก_4_บุคลากร_CM_W{week_no}.xlsx", fb)
        except Exception as e:
            zf.writestr(f"ERROR_appendix4.txt", f"Error: {e}")

    return zip_buf.getvalue()
